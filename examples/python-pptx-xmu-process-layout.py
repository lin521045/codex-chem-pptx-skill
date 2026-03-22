#!/usr/bin/env python3
"""Generate a process-design deck with richer XMU-style visuals."""

from __future__ import annotations

import re
from dataclasses import dataclass

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


@dataclass
class SlideSpec:
    chapter_name: str
    slide_title: str
    bullet_points: list[str]


XMU_BLUE = RGBColor(0x00, 0x3F, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF8, 0xF9, 0xFA)
TEXT = RGBColor(0x33, 0x33, 0x33)
DIM = RGBColor(0x88, 0x88, 0x88)
SIDEBAR_BG = RGBColor(0xF1, 0xF4, 0xF8)
LINE = RGBColor(0xD9, 0xE2, 0xEC)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
PALE_BLUE = RGBColor(0xEE, 0xF4, 0xFB)
PALE_GREEN = RGBColor(0xEC, 0xF8, 0xF2)
PALE_GOLD = RGBColor(0xFF, 0xF6, 0xE8)
PALE_RED = RGBColor(0xFB, 0xEE, 0xEE)
GREEN = RGBColor(0x12, 0x7C, 0x5A)
ORANGE = RGBColor(0xD9, 0x77, 0x06)
RED = RGBColor(0xC2, 0x41, 0x0C)

SLIDE_W = 13.333
SLIDE_H = 7.5
BANNER_H = 0.9
SIDEBAR_W = 2.7
SAFE_X = 3.35
SAFE_Y = 1.5
SAFE_W = 9.35
SAFE_H = 5.6


def insert_cn_en_spacing(text: str) -> str:
    text = re.sub(r"([\u4e00-\u9fff])([A-Za-z0-9])", r"\1 \2", text)
    text = re.sub(r"([A-Za-z0-9])([\u4e00-\u9fff])", r"\1 \2", text)
    return text


def normalize_units(text: str) -> str:
    replacements = {
        "mol / L": "mol/L",
        "kJ / mol": "kJ/mol",
        "m3/h": "m³/h",
        "m3 / h": "m³/h",
        "wt %": "wt%",
        "° C": "°C",
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    text = re.sub(r"(\d)(mol|kg|t|MPa|bar|°C|h|wt%|kmol/h|t/h|m³/h)\b", r"\1 \2", text)
    return text


def normalize_chem_text(text: str) -> str:
    text = insert_cn_en_spacing(text)
    text = normalize_units(text)
    text = text.replace("CO2", "CO₂").replace("H2O", "H₂O")
    return text


def add_banner(slide, title: str) -> None:
    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(SLIDE_W),
        Inches(BANNER_H),
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = XMU_BLUE
    banner.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.14), Inches(SLIDE_W - 1.1), Inches(0.56))
    frame = title_box.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = normalize_chem_text(title)
    run.font.name = "Arial"
    run.font.size = Pt(23)
    run.font.bold = True
    run.font.color.rgb = WHITE


def add_sidebar(slide, chapters: list[str], active_index: int) -> None:
    sidebar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(BANNER_H),
        Inches(SIDEBAR_W),
        Inches(SLIDE_H - BANNER_H),
    )
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = SIDEBAR_BG
    sidebar.line.fill.background()

    label_box = slide.shapes.add_textbox(Inches(0.28), Inches(1.08), Inches(2.0), Inches(0.25))
    label_frame = label_box.text_frame
    label_frame.clear()
    p = label_frame.paragraphs[0]
    run = p.add_run()
    run.text = "PROCESS NAVIGATION"
    run.font.name = "Arial"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = DIM

    nav_top = 1.45
    nav_step = 0.82
    for idx, chapter in enumerate(chapters):
        y = nav_top + idx * nav_step
        active = idx == active_index
        if active:
            highlight = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(0.18),
                Inches(y - 0.05),
                Inches(2.28),
                Inches(0.52),
            )
            highlight.fill.solid()
            highlight.fill.fore_color.rgb = XMU_BLUE
            highlight.line.fill.background()

        box = slide.shapes.add_textbox(Inches(0.34), Inches(y), Inches(2.0), Inches(0.34))
        frame = box.text_frame
        frame.clear()
        p = frame.paragraphs[0]
        run = p.add_run()
        run.text = chapter
        run.font.name = "Arial"
        run.font.size = Pt(13 if active else 11.5)
        run.font.bold = active
        run.font.color.rgb = WHITE if active else DIM


def add_safe_zone(slide) -> None:
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(SAFE_X),
        Inches(SAFE_Y),
        Inches(SAFE_W),
        Inches(SAFE_H),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = LINE
    box.line.width = Pt(1)


def add_bullets(slide, bullets: list[str]) -> None:
    body = slide.shapes.add_textbox(Inches(SAFE_X + 0.3), Inches(SAFE_Y + 0.32), Inches(4.7), Inches(2.05))
    frame = body.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP

    for idx, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.bullet = True
        p.space_after = Pt(7)
        run = p.add_run()
        run.text = normalize_chem_text(bullet)
        run.font.name = "Arial"
        run.font.size = Pt(15)
        run.font.color.rgb = TEXT


def add_card(slide, x, y, w, h, title, fill=WHITE, title_color=XMU_BLUE):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = LINE
    card.line.width = Pt(1)

    title_box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.12), Inches(w - 0.36), Inches(0.26))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Arial"
    run.font.size = Pt(12.5)
    run.font.bold = True
    run.font.color.rgb = title_color
    return card


def add_metric_chip(slide, x, y, w, label, value, fill):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.78))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.fill.background()

    t1 = slide.shapes.add_textbox(Inches(x + 0.12), Inches(y + 0.1), Inches(w - 0.24), Inches(0.18))
    p1 = t1.text_frame.paragraphs[0]
    r1 = p1.add_run()
    r1.text = label
    r1.font.name = "Arial"
    r1.font.size = Pt(9.5)
    r1.font.bold = True
    r1.font.color.rgb = DIM

    t2 = slide.shapes.add_textbox(Inches(x + 0.12), Inches(y + 0.3), Inches(w - 0.24), Inches(0.25))
    p2 = t2.text_frame.paragraphs[0]
    r2 = p2.add_run()
    r2.text = value
    r2.font.name = "Arial"
    r2.font.size = Pt(16)
    r2.font.bold = True
    r2.font.color.rgb = XMU_BLUE


def add_text(slide, x, y, w, h, text, size=11, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = normalize_chem_text(text)
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_arrow(slide, x1, y1, x2, y2, label=None):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x1), Inches(y1), Inches(x2 - x1), Inches(y2 - y1))
    line.fill.solid()
    line.fill.fore_color.rgb = XMU_BLUE
    line.line.fill.background()
    if label:
        add_text(slide, x1, y1 - 0.2, x2 - x1, 0.18, label, size=8.5, color=DIM, align=PP_ALIGN.CENTER)


def render_overview(slide):
    add_metric_chip(slide, 8.2, 1.86, 1.35, "Design Capacity", "10 万吨/年", PALE_BLUE)
    add_metric_chip(slide, 9.7, 1.86, 1.15, "Operating", "8000 h/a", PALE_GREEN)
    add_metric_chip(slide, 11.0, 1.86, 1.1, "Route", "Non-Phosgene", PALE_GOLD)

    add_card(slide, 8.15, 2.95, 4.22, 2.35, "项目边界与产品定位", fill=PALE_BLUE)
    add_text(slide, 8.35, 3.28, 3.8, 0.28, "Feed System", size=10, bold=True)
    add_text(slide, 8.35, 3.62, 3.8, 0.52, "BPA、DPC 与公用工程系统\n与 PC 聚合主体装置进行热集成。", size=10.5)
    add_text(slide, 8.35, 4.38, 3.8, 0.28, "Product Slate", size=10, bold=True)
    add_text(slide, 8.35, 4.72, 3.8, 0.48, "目标产品为光学级与工程塑料级 Polycarbonate (PC)，并兼顾副产 Phenol 回收。", size=10.5)


def render_route_screening(slide):
    add_card(slide, 8.0, 1.95, 1.1, 1.0, "原料\n预处理", fill=PALE_BLUE)
    add_card(slide, 9.4, 1.95, 1.1, 1.0, "酯交换\n反应", fill=PALE_GREEN)
    add_card(slide, 10.8, 1.95, 1.1, 1.0, "脱挥\n精制", fill=PALE_GOLD)
    add_card(slide, 12.2, 1.95, 0.7, 1.0, "PC", fill=PALE_BLUE)
    add_arrow(slide, 9.1, 2.23, 9.36, 2.48, "F-101")
    add_arrow(slide, 10.5, 2.23, 10.76, 2.48, "F-201")
    add_arrow(slide, 11.9, 2.23, 12.16, 2.48, "F-301")

    add_card(slide, 8.0, 3.45, 4.9, 1.85, "推荐流程骨架", fill=WHITE)
    add_text(slide, 8.22, 3.84, 4.45, 0.42, "上图不再是通用占位框，而是类似 Aspen PFD 的流程链骨架，可继续替换为真实流程图。", size=10.5)
    add_text(slide, 8.22, 4.48, 4.45, 0.42, "建议在正式版本中接入流股编号、主要温压条件和 Phenol 回收回路。", size=10.5)


def render_balance(slide):
    add_card(slide, 8.0, 1.94, 4.9, 3.35, "关键流股与热负荷", fill=WHITE)

    table = slide.shapes.add_table(5, 3, Inches(8.18), Inches(2.26), Inches(2.75), Inches(2.55)).table
    headers = ["流股", "流量", "说明"]
    rows = [
        ("BPA Feed", "9.8 t/h", "主原料"),
        ("DPC Feed", "10.6 t/h", "碳酸酯供给"),
        ("Phenol Recycle", "4.2 t/h", "回收回路"),
        ("PC Product", "12.5 t/h", "目标产出"),
    ]
    for col, text in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = XMU_BLUE
        para = cell.text_frame.paragraphs[0]
        para.runs[0].font.bold = True
        para.runs[0].font.color.rgb = WHITE
        para.runs[0].font.size = Pt(10)
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_idx % 2 else BG
            para = cell.text_frame.paragraphs[0]
            para.runs[0].font.size = Pt(9.5)
            para.runs[0].font.color.rgb = TEXT

    add_text(slide, 11.25, 2.2, 1.35, 0.22, "Heat Duty", size=10, bold=True, color=XMU_BLUE)
    duties = [("反应段", 0.95, GREEN), ("脱挥段", 0.72, ORANGE), ("塔系", 1.18, RED)]
    bar_y = 2.62
    for name, width, color in duties:
        slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(11.25), Inches(bar_y), Inches(width), Inches(0.22)).fill.solid()
        slide.shapes[-1].fill.fore_color.rgb = color
        slide.shapes[-1].line.fill.background()
        add_text(slide, 11.25, bar_y - 0.18, 1.4, 0.16, name, size=8.5, color=DIM)
        bar_y += 0.72


def render_equipment(slide):
    add_card(slide, 8.0, 1.95, 4.9, 3.35, "设备清单与功能分区", fill=WHITE)

    table = slide.shapes.add_table(4, 4, Inches(8.16), Inches(2.22), Inches(4.48), Inches(1.85)).table
    headers = ["设备位号", "设备名称", "材质", "关键点"]
    rows = [
        ("R-101", "酯交换反应器", "SS316L", "高黏体系传质"),
        ("V-201", "脱挥器", "SS316L", "真空稳定性"),
        ("T-301", "Phenol 塔", "2205", "腐蚀与回收"),
    ]
    for col, text in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = XMU_BLUE
        para = cell.text_frame.paragraphs[0]
        para.runs[0].font.bold = True
        para.runs[0].font.color.rgb = WHITE
        para.runs[0].font.size = Pt(9.5)
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_idx % 2 else BG
            para = cell.text_frame.paragraphs[0]
            para.runs[0].font.size = Pt(9)
            para.runs[0].font.color.rgb = TEXT

    add_card(slide, 8.18, 4.28, 1.18, 0.78, "R-101", fill=PALE_BLUE)
    add_card(slide, 9.72, 4.28, 1.18, 0.78, "V-201", fill=PALE_GREEN)
    add_card(slide, 11.26, 4.28, 1.18, 0.78, "T-301", fill=PALE_GOLD)
    add_arrow(slide, 9.36, 4.54, 9.68, 4.74)
    add_arrow(slide, 10.9, 4.54, 11.22, 4.74)


def render_hse_tea(slide):
    add_card(slide, 8.0, 1.94, 2.45, 3.35, "Risk Matrix", fill=WHITE)
    labels_x = 8.24
    labels_y = 2.26
    cell_w = 0.54
    cell_h = 0.46
    colors = [
        [PALE_GREEN, PALE_GREEN, PALE_GOLD],
        [PALE_GREEN, PALE_GOLD, PALE_RED],
        [PALE_GOLD, PALE_RED, PALE_RED],
    ]
    for r in range(3):
        for c in range(3):
            cell = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(labels_x + c * cell_w),
                Inches(labels_y + r * cell_h),
                Inches(cell_w - 0.03),
                Inches(cell_h - 0.03),
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = colors[r][c]
            cell.line.color.rgb = LINE
    add_text(slide, 8.22, 3.82, 1.95, 0.5, "重点关注高温熔体泄漏、真空失效和 Phenol 回收系统波动。", size=9.5)

    add_card(slide, 10.7, 1.94, 2.2, 1.5, "CAPEX / OPEX 摘要", fill=WHITE)
    add_metric_chip(slide, 10.92, 2.28, 0.82, "CAPEX", "Base", PALE_BLUE)
    add_metric_chip(slide, 11.84, 2.28, 0.82, "OPEX", "Optimized", PALE_GREEN)
    add_text(slide, 10.9, 3.0, 1.8, 0.22, "Phenol 回收率和热集成效率是关键敏感项。", size=9.5)

    add_card(slide, 10.7, 3.68, 2.2, 1.62, "建议动作", fill=WHITE)
    add_text(slide, 10.92, 4.02, 1.78, 0.78, "1. 完成 HAZOP / LOPA\n2. 细化塔系热集成\n3. 更新 CAPEX / OPEX 模型", size=9.5)


def create_content_slide(prs: Presentation, chapters: list[str], active_index: int, spec: SlideSpec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_banner(slide, spec.slide_title)
    add_sidebar(slide, chapters, active_index)
    add_safe_zone(slide)
    add_bullets(slide, spec.bullet_points)

    if active_index == 0:
        render_overview(slide)
    elif active_index == 1:
        render_route_screening(slide)
    elif active_index == 2:
        render_balance(slide)
    elif active_index == 3:
        render_equipment(slide)
    else:
        render_hse_tea(slide)
    return slide


def build_process_specs() -> list[SlideSpec]:
    return [
        SlideSpec(
            chapter_name="项目背景与产能规划",
            slide_title="项目背景与产能规划",
            bullet_points=[
                "目标产品为 Polycarbonate (PC)，设计产能为 10 万吨/年",
                "路线选择需兼顾高分子量控制、Phenol 回收与 HSE 风险",
                "装置按 8000 h/a 运行，并预留与上游 DPC 系统热集成接口",
            ],
        ),
        SlideSpec(
            chapter_name="工艺路线比选",
            slide_title="工艺路线比选与推荐方案",
            bullet_points=[
                "比较光气法与非光气熔融酯交换法的安全性、成熟度与环保压力",
                "非光气路线更符合当前政策环境，并具备更好的副产 Phenol 回收潜力",
                "推荐采用 DPC 与 BPA 熔融酯交换路线作为基准方案",
            ],
        ),
        SlideSpec(
            chapter_name="物料与能量衡算",
            slide_title="物料衡算 (Mass Balance) 与能量衡算 (Energy Balance)",
            bullet_points=[
                "PC 目标产出约 12.5 t/h，BPA 进料约 9.8 t/h，DPC 进料约 10.6 t/h",
                "回收 Phenol 流量约 4.2 t/h，需重点考虑精馏系统负荷",
                "主要热负荷集中在反应段、脱挥段与塔系再沸器，可通过热油系统集成",
            ],
        ),
        SlideSpec(
            chapter_name="关键设备选型",
            slide_title="关键设备选型",
            bullet_points=[
                "反应器需适配高黏熔体体系，并强化传质与停留时间分布控制",
                "脱挥器与精馏塔材质需兼顾高温与 Phenol 腐蚀环境",
                "关键转动设备按 1 用 1 备配置，提升装置连续运行稳定性",
            ],
        ),
        SlideSpec(
            chapter_name="HSE 与技术经济评价",
            slide_title="HSE 与技术经济评价 (Techno-Economic Analysis)",
            bullet_points=[
                "非光气路线显著降低光气相关重大毒害风险，但需关注高温熔体泄漏与真空失效",
                "建议对反应段、脱挥段和精馏塔系统开展 HAZOP / LOPA 分析",
                "在高 Phenol 回收率与公用工程优化条件下，项目具备良好经济可行性",
            ],
        ),
    ]


def main() -> int:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    specs = build_process_specs()
    chapters = [item.chapter_name for item in specs]

    for index, spec in enumerate(specs):
        create_content_slide(prs, chapters, index, spec)

    prs.save("xmu-chem-process-design-sample.pptx")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
