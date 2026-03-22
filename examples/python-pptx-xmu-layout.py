#!/usr/bin/env python3
"""Generate a chemistry academic-defense deck with a top banner and left sidebar."""

from __future__ import annotations

import re
from dataclasses import dataclass

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


@dataclass
class SlideSpec:
    chapter_name: str
    slide_title: str
    bullet_points: list[str]
    visual_placeholder: str


XMU_BLUE = RGBColor(0x00, 0x3F, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF8, 0xF9, 0xFA)
TEXT = RGBColor(0x33, 0x33, 0x33)
DIM = RGBColor(0x88, 0x88, 0x88)
SIDEBAR_BG = RGBColor(0xF1, 0xF4, 0xF8)
LINE = RGBColor(0xD9, 0xE2, 0xEC)
PLACEHOLDER_BG = RGBColor(0xEE, 0xF4, 0xFB)

SLIDE_W = 13.333
SLIDE_H = 7.5
BANNER_H = 0.9
SIDEBAR_W = 2.7
SAFE_X = 3.35
SAFE_Y = 1.5
SAFE_W = 9.35
SAFE_H = 5.6

SUBSCRIPT_MAP = str.maketrans("0123456789+-()", "₀₁₂₃₄₅₆₇₈₉₊₋₍₎")
SUPERSCRIPT_MAP = str.maketrans("0123456789+-()", "⁰¹²³⁴⁵⁶⁷⁸⁹⁺⁻⁽⁾")


def insert_cn_en_spacing(text: str) -> str:
    text = re.sub(r"([\u4e00-\u9fff])([A-Za-z0-9])", r"\1 \2", text)
    text = re.sub(r"([A-Za-z0-9])([\u4e00-\u9fff])", r"\1 \2", text)
    return text


def normalize_units(text: str) -> str:
    replacements = {
        "mol / L": "mol/L",
        "kJ / mol": "kJ/mol",
        "m3/h": "m³/h",
        "m3 / h": "m³/h",
        "wt %": "wt%",
        "° C": "°C",
        "cm-1": "cm⁻¹",
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    text = re.sub(r"(\d)(mmol|mol|mA|A|V|h|wt%|nm|eV|cm⁻¹|°C|MPa)\b", r"\1 \2", text)
    return text


def to_subscript(value: str) -> str:
    return value.translate(SUBSCRIPT_MAP)


def to_superscript(value: str) -> str:
    return value.translate(SUPERSCRIPT_MAP)


def normalize_chem_text(text: str) -> str:
    text = insert_cn_en_spacing(text)
    text = normalize_units(text)
    text = re.sub(r"\b(\d{1,3})([A-Z][a-z]?)(?=\s*(NMR|MS))", lambda m: f"{to_superscript(m.group(1))}{m.group(2)}", text)
    text = re.sub(r"([A-Za-z\)])(\d+)", lambda m: m.group(1) + to_subscript(m.group(2)), text)
    text = re.sub(r"\^([0-9+-]+)", lambda m: to_superscript(m.group(1)), text)
    return text


def add_banner(slide, title: str) -> None:
    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(SLIDE_W),
        Inches(BANNER_H),
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = XMU_BLUE
    banner.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.14), Inches(SLIDE_W - 1.1), Inches(0.56))
    frame = title_box.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = normalize_chem_text(title)
    run.font.name = "Arial"
    run.font.size = Pt(23)
    run.font.bold = True
    run.font.color.rgb = WHITE


def add_sidebar(slide, chapters: list[str], active_index: int) -> None:
    sidebar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(BANNER_H),
        Inches(SIDEBAR_W),
        Inches(SLIDE_H - BANNER_H),
    )
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = SIDEBAR_BG
    sidebar.line.fill.background()

    label_box = slide.shapes.add_textbox(Inches(0.28), Inches(1.08), Inches(2.0), Inches(0.25))
    label_frame = label_box.text_frame
    label_frame.clear()
    p = label_frame.paragraphs[0]
    run = p.add_run()
    run.text = "GLOBAL NAVIGATION"
    run.font.name = "Arial"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = DIM

    nav_top = 1.45
    nav_step = 0.82
    for idx, chapter in enumerate(chapters):
        y = nav_top + idx * nav_step
        active = idx == active_index
        if active:
            highlight = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(0.18),
                Inches(y - 0.05),
                Inches(2.28),
                Inches(0.52),
            )
            highlight.fill.solid()
            highlight.fill.fore_color.rgb = XMU_BLUE
            highlight.line.fill.background()

        box = slide.shapes.add_textbox(Inches(0.34), Inches(y), Inches(2.0), Inches(0.34))
        frame = box.text_frame
        frame.clear()
        p = frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = chapter
        run.font.name = "Arial"
        run.font.size = Pt(14 if active else 12)
        run.font.bold = active
        run.font.color.rgb = WHITE if active else DIM


def add_safe_zone(slide) -> None:
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(SAFE_X),
        Inches(SAFE_Y),
        Inches(SAFE_W),
        Inches(SAFE_H),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = LINE
    box.line.width = Pt(1)


def add_body_and_placeholder(slide, spec: SlideSpec) -> None:
    left_x = SAFE_X + 0.35
    top_y = SAFE_Y + 0.35
    content_w = SAFE_W * 0.56
    placeholder_x = SAFE_X + SAFE_W * 0.63
    placeholder_w = SAFE_W * 0.30

    body_box = slide.shapes.add_textbox(Inches(left_x), Inches(top_y), Inches(content_w), Inches(4.6))
    frame = body_box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP

    for idx, bullet in enumerate(spec.bullet_points):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        p.bullet = True
        add_chem_runs(p, normalize_chem_text(bullet), Pt(16))

    placeholder = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(placeholder_x),
        Inches(top_y),
        Inches(placeholder_w),
        Inches(4.2),
    )
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = PLACEHOLDER_BG
    placeholder.line.color.rgb = XMU_BLUE
    placeholder.line.width = Pt(1.5)

    placeholder_box = slide.shapes.add_textbox(
        Inches(placeholder_x + 0.18),
        Inches(top_y + 0.2),
        Inches(placeholder_w - 0.36),
        Inches(3.8),
    )
    placeholder_frame = placeholder_box.text_frame
    placeholder_frame.clear()
    placeholder_frame.word_wrap = True
    placeholder_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = placeholder_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = normalize_chem_text(spec.visual_placeholder)
    run.font.name = "Arial"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = XMU_BLUE


def add_chem_runs(paragraph, text: str, font_size: Pt) -> None:
    """Use run-level formatting when available; fall back to pre-normalized Unicode text."""
    formula_match = re.fullmatch(r"(.*?)([A-Z][a-z]?)(\d+)(.*)", text)
    if formula_match:
        prefix, element, digits, suffix = formula_match.groups()
        probe = paragraph.add_run()
        if hasattr(probe.font, "subscript"):
            if prefix:
                probe.text = prefix
                style_run(probe, font_size)
            else:
                probe.text = element
                style_run(probe, font_size)
                element = ""

            if element:
                run = paragraph.add_run()
                run.text = element
                style_run(run, font_size)
            run = paragraph.add_run()
            run.text = digits
            style_run(run, Pt(font_size.pt - 3))
            run.font.subscript = True
            if suffix:
                run = paragraph.add_run()
                run.text = suffix
                style_run(run, font_size)
            return
        paragraph._p.remove(probe._r)

    run = paragraph.add_run()
    run.text = text
    style_run(run, font_size)


def style_run(run, size: Pt) -> None:
    run.font.name = "Arial"
    run.font.size = size
    run.font.color.rgb = TEXT


def create_content_slide(prs: Presentation, chapters: list[str], active_index: int, spec: SlideSpec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_banner(slide, spec.slide_title)
    add_sidebar(slide, chapters, active_index)
    add_safe_zone(slide)
    add_body_and_placeholder(slide, spec)
    return slide


def build_demo_specs() -> list[SlideSpec]:
    return [
        SlideSpec(
            chapter_name="研究背景与意义",
            slide_title="研究背景与意义",
            bullet_points=[
                "CO₂ 电还原可实现 C1 化学品的低碳转化",
                "Fe-N-C 兼具成本优势与较高 FECO 潜力",
                "本研究关注 MOF 衍生 Fe-N₄ 位点的结构 - 性能关系",
            ],
            visual_placeholder="[在此插入 CO₂RR 文献综述图或 ChemDraw 结构图]",
        ),
        SlideSpec(
            chapter_name="Reaction Mechanism",
            slide_title="Reaction Mechanism 假设与位点设计",
            bullet_points=[
                "以 ZIF-8 为前驱体限域 Fe 物种，热解后形成分散 Fe-N₄ 位点",
                "假设 Fe-N₄ 可优先稳定 *COOH 中间体并促进 CO 生成",
                "结合 In-situ FTIR 与 DFT 分析关键中间体演变",
            ],
            visual_placeholder="[在此插入 Reaction Mechanism 示意图]",
        ),
        SlideSpec(
            chapter_name="实验方法与表征",
            slide_title="实验方法与表征",
            bullet_points=[
                "采用 H-type 电解池，在 0.5 mol/L KHCO₃ 中测试 CO₂RR 活性",
                "通过 XRD、XPS、HAADF-STEM 与 XANES 确认 Fe-N₄ 结构",
                "产物由 GC 和 ¹H NMR 联合定量，确保碳平衡闭合",
            ],
            visual_placeholder="[在此插入 XRD / XPS / SEM / TEM 表征总览图]",
        ),
        SlideSpec(
            chapter_name="结果与讨论",
            slide_title="结果与讨论",
            bullet_points=[
                "在 -0.72 V vs. RHE 下 FECO 达到 94%，jCO 达 28.6 mA/cm²",
                "HAADF-STEM 未观察到明显 Fe 团簇，支持单原子分散特征",
                "In-situ FTIR 捕获 *COOH 与 *CO 信号，支持所提机制",
            ],
            visual_placeholder="[在此插入 FECO / jCO 曲线与原位谱图]",
        ),
        SlideSpec(
            chapter_name="结论与展望",
            slide_title="结论与展望",
            bullet_points=[
                "MOF 衍生策略可构建高分散 Fe-N₄ 位点并提升 CO₂RR 性能",
                "多尺度表征支持位点结构与高 CO 选择性之间的关联",
                "下一步将拓展到 MEA 条件，并引入 operando XAS 验证",
            ],
            visual_placeholder="[在此插入未来工作路线图或 MEA 平台示意图]",
        ),
    ]


def main() -> int:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    specs = build_demo_specs()
    chapters = [item.chapter_name for item in specs]

    for index, spec in enumerate(specs):
        create_content_slide(prs, chapters, index, spec)

    prs.save("xmu-chem-academic-defense-sample.pptx")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
