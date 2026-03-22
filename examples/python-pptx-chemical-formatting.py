from pptx import Presentation
from pptx.util import Inches, Pt


def add_formula_runs(paragraph):
    run1 = paragraph.add_run()
    run1.text = "SO"
    run1.font.size = Pt(20)

    run2 = paragraph.add_run()
    run2.text = "4"
    run2.font.size = Pt(14)
    if hasattr(run2.font, "subscript"):
        run2.font.subscript = True

    run3 = paragraph.add_run()
    run3.text = "2-"
    run3.font.size = Pt(12)
    if hasattr(run3.font, "superscript"):
        run3.font.superscript = True


prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
textbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(8.0), Inches(1.2))
tf = textbox.text_frame
tf.clear()

p = tf.paragraphs[0]
title_run = p.add_run()
title_run.text = "硫酸根示例："
title_run.font.size = Pt(20)
add_formula_runs(p)

textbox2 = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(10.0), Inches(1.2))
tf2 = textbox2.text_frame
tf2.text = "如果运行级上标/下标不可用，请直接输出 Unicode：SO₄²⁻、H₂O、¹³C NMR。"

prs.save("python-pptx-chemical-formatting-demo.pptx")
