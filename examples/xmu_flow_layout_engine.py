#!/usr/bin/env python3
"""Reusable XMU-style flow layout engine for chemistry decks."""

from __future__ import annotations

import json
import math
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


XMU_BLUE = RGBColor(0x00, 0x3F, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF8, 0xF9, 0xFA)
TEXT = RGBColor(0x33, 0x33, 0x33)
DIM = RGBColor(0x88, 0x88, 0x88)
LINE = RGBColor(0xD9, 0xE2, 0xEC)
SIDEBAR_BG = RGBColor(0xF1, 0xF4, 0xF8)
CARD_BG = RGBColor(0xFB, 0xFC, 0xFD)
PLACEHOLDER_BG = RGBColor(0xEE, 0xF4, 0xFB)
CALLOUT_BG = RGBColor(0xF1, 0xF3, 0xF5)
GREEN = RGBColor(0x12, 0x7C, 0x5A)
ORANGE = RGBColor(0xD9, 0x77, 0x06)
RED = RGBColor(0xC2, 0x41, 0x0C)

SLIDE_W = 13.333
SLIDE_H = 7.5
BANNER_H = 0.86
SIDEBAR_W = 2.58
SAFE_X = 2.92
SAFE_Y = 1.12
SAFE_W = 9.95
SAFE_H = 6.02

SUBSCRIPT_MAP = str.maketrans("0123456789+-()", "₀₁₂₃₄₅₆₇₈₉₊₋₍₎")
SUPERSCRIPT_MAP = str.maketrans("0123456789+-()", "⁰¹²³⁴⁵⁶⁷⁸⁹⁺⁻⁽⁾")
CJK = r"\u4e00-\u9fff"
ASCII_BLOCK = r"A-Za-z0-9"
FORMULA_RE = re.compile(r"(?<![A-Za-z])((?:\d+)?(?:[A-Z][a-z]?\d*)+(?:\^[0-9]*[+-]|[+-])?)(?![A-Za-z])")
ISOTOPE_RE = re.compile(r"\b(\d{1,3})([A-Z][a-z]?)(?=\s*(?:NMR|MS|label|标记))")
ALLOWED_PLACEHOLDERS = {"None", "ChemDraw", "Aspen PFD", "Chart", "Characterization", "Equipment", "Risk Matrix"}


@dataclass
class BulletPoint:
    keyword: str
    description: str


@dataclass
class SlideSpec:
    chapter: str
    slide_title: str
    bullet_points: list[BulletPoint]
    key_takeaway: str
    placeholder_type: str = "None"


def insert_cn_en_spacing(text: str) -> str:
    text = re.sub(fr"([{CJK}])([{ASCII_BLOCK}])", r"\1 \2", text)
    text = re.sub(fr"([{ASCII_BLOCK}])([{CJK}])", r"\1 \2", text)
    return re.sub(r"\s{2,}", " ", text).strip()


def normalize_units(text: str) -> str:
    replacements = {
        "mol / L": "mol/L",
        "kJ / mol": "kJ/mol",
        "kg / h": "kg/h",
        "t / h": "t/h",
        "m3/h": "m³/h",
        "m3 / h": "m³/h",
        "Nm3/h": "Nm³/h",
        "wt %": "wt%",
        "vol %": "vol%",
        "° C": "°C",
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    return re.sub(r"(\d)(mmol|mol|g|kg|t|wt%|MPa|kPa|bar|°C|K|h|min|s|t/h|kg/h|m³/h)\b", r"\1 \2", text)


def to_subscript(value: str) -> str:
    return value.translate(SUBSCRIPT_MAP)


def to_superscript(value: str) -> str:
    return value.translate(SUPERSCRIPT_MAP)


def format_formula_token(token: str) -> str:
    isotope = ""
    core = token
    match = re.match(r"^(\d+)([A-Z].*)$", token)
    if match and re.search(r"\d", match.group(2)):
        isotope = to_superscript(match.group(1))
        core = match.group(2)

    charge = ""
    if "^" in core:
        core, charge = core.split("^", 1)
    elif core.endswith(("+", "-")) and re.search(r"\d", core):
        core, charge = core[:-1], core[-1]

    core = re.sub(r"([A-Za-z\)])(\d+)", lambda m: m.group(1) + to_subscript(m.group(2)), core)
    return isotope + core + (to_superscript(charge) if charge else "")


def normalize_chemistry_unicode(text: str) -> str:
    text = ISOTOPE_RE.sub(lambda m: f"{to_superscript(m.group(1))}{m.group(2)}", text)

    def repl(match: re.Match[str]) -> str:
        token = match.group(1)
        if not re.search(r"[\d\^+-]", token):
            return token
        return format_formula_token(token)

    return FORMULA_RE.sub(repl, text)


def normalize_text(text: str) -> str:
    return normalize_chemistry_unicode(normalize_units(insert_cn_en_spacing(text or "")))


def display_units(text: str) -> float:
    total = 0.0
    for char in text:
        if re.match(fr"[{CJK}]", char):
            total += 1.0
        elif char.isspace():
            total += 0.35
        elif char in "ilI.,:;|![]()/":
            total += 0.35
        elif ord(char) > 127:
            total += 0.75
        else:
            total += 0.62
    return total


def shorten_text(text: str, max_units: float) -> str:
    text = normalize_text(text)
    if display_units(text) <= max_units:
        return text
    running = 0.0
    chars: list[str] = []
    for char in text:
        char_units = display_units(char)
        if running + char_units > max_units - 1:
            break
        chars.append(char)
        running += char_units
    return "".join(chars).rstrip(" ，,;；:：") + "…"


def estimate_line_capacity(width_in: float, font_size_pt: float) -> float:
    return max(18.0, width_in * 7.8 * (12.0 / font_size_pt))


def estimate_line_count(text: str, width_in: float, font_size_pt: float) -> int:
    return max(1, math.ceil(display_units(text) / estimate_line_capacity(width_in, font_size_pt)))


def estimate_textbox_height(text: str, width_in: float, font_size_pt: float, padding_in: float = 0.08) -> float:
    lines = estimate_line_count(text, width_in, font_size_pt)
    line_height = (font_size_pt / 72.0) * 1.35
    return max(0.24, lines * line_height + padding_in * 2)


def bullet_from_dict(item: dict[str, Any]) -> BulletPoint:
    return BulletPoint(
        keyword=shorten_text(str(item.get("keyword", "要点")), 24),
        description=shorten_text(str(item.get("description", "")), 60),
    )


def slide_from_dict(item: dict[str, Any]) -> SlideSpec:
    placeholder_type = str(item.get("placeholder_type", "None"))
    if placeholder_type not in ALLOWED_PLACEHOLDERS:
        placeholder_type = "None"
    bullets = [bullet_from_dict(point) for point in item.get("bullet_points", [])][:4]
    return SlideSpec(
        chapter=shorten_text(str(item.get("chapter", "未命名章节")), 26),
        slide_title=shorten_text(str(item.get("slide_title", "未命名页面")), 40),
        bullet_points=bullets,
        key_takeaway=shorten_text(str(item.get("key_takeaway", "请补充关键结论。")), 72),
        placeholder_type=placeholder_type,
    )


def split_dense_slides(slides: list[SlideSpec], max_bullets: int = 4) -> list[SlideSpec]:
    expanded: list[SlideSpec] = []
    for slide in slides:
        if len(slide.bullet_points) <= max_bullets:
            expanded.append(slide)
            continue
        chunks = [slide.bullet_points[i : i + max_bullets] for i in range(0, len(slide.bullet_points), max_bullets)]
        for index, chunk in enumerate(chunks, start=1):
            suffix = f" - Part {index}" if len(chunks) > 1 else ""
            expanded.append(
                SlideSpec(
                    chapter=slide.chapter,
                    slide_title=f"{slide.slide_title}{suffix}",
                    bullet_points=chunk,
                    key_takeaway=slide.key_takeaway,
                    placeholder_type=slide.placeholder_type,
                )
            )
    return expanded


def load_slide_specs(payload: str | Path | list[dict[str, Any]]) -> list[SlideSpec]:
    if isinstance(payload, list):
        data = payload
    elif isinstance(payload, Path):
        data = json.loads(payload.read_text(encoding="utf-8"))
    else:
        raw = str(payload).strip()
        if raw.startswith("[") or raw.startswith("{"):
            data = json.loads(raw)
        else:
            data = json.loads(Path(raw).read_text(encoding="utf-8"))
    if isinstance(data, dict):
        data = data.get("slides", [])
    slides = [slide_from_dict(item) for item in data]
    return split_dense_slides(slides)


def chapter_outline(slides: list[SlideSpec]) -> list[str]:
    ordered: list[str] = []
    for slide in slides:
        if slide.chapter not in ordered:
            ordered.append(slide.chapter)
    return ordered


def style_shape_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def apply_text_frame_style(text_frame, vertical_anchor=MSO_ANCHOR.TOP) -> None:
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.vertical_anchor = vertical_anchor
    text_frame.margin_left = Inches(0.04)
    text_frame.margin_right = Inches(0.04)
    text_frame.margin_top = Inches(0.02)
    text_frame.margin_bottom = Inches(0.02)


def add_textbox(slide, x: float, y: float, w: float, h: float, vertical_anchor=MSO_ANCHOR.TOP):
    textbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    apply_text_frame_style(textbox.text_frame, vertical_anchor=vertical_anchor)
    textbox.text_frame.clear()
    return textbox


def add_text_run(paragraph, text: str, size_pt: float, color: RGBColor, bold: bool = False) -> None:
    run = paragraph.add_run()
    run.text = normalize_text(text)
    run.font.name = "Arial"
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def add_banner(slide, title: str) -> None:
    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(SLIDE_W),
        Inches(BANNER_H),
    )
    style_shape_fill(banner, XMU_BLUE)
    banner.line.fill.background()

    title_box = add_textbox(slide, 0.65, 0.12, SLIDE_W - 1.3, 0.54, vertical_anchor=MSO_ANCHOR.MIDDLE)
    paragraph = title_box.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    add_text_run(paragraph, title, 23, WHITE, bold=True)


def add_sidebar(slide, chapters: list[str], active_chapter: str) -> None:
    sidebar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(BANNER_H),
        Inches(SIDEBAR_W),
        Inches(SLIDE_H - BANNER_H),
    )
    style_shape_fill(sidebar, SIDEBAR_BG)
    sidebar.line.fill.background()

    label_box = add_textbox(slide, 0.28, 1.0, 1.95, 0.26)
    paragraph = label_box.text_frame.paragraphs[0]
    add_text_run(paragraph, "全局导航 | OUTLINE", 9.2, DIM, bold=True)

    nav_start_y = 1.42
    nav_available = SLIDE_H - nav_start_y - 0.55
    item_pitch = nav_available / max(1, len(chapters))
    item_height = min(0.56, max(0.42, item_pitch - 0.08))

    for index, chapter in enumerate(chapters):
        y = nav_start_y + index * item_pitch
        is_active = chapter == active_chapter
        if is_active:
            highlight = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(0.18),
                Inches(y - 0.03),
                Inches(SIDEBAR_W - 0.34),
                Inches(item_height),
            )
            style_shape_fill(highlight, XMU_BLUE)
            highlight.line.fill.background()

        box = add_textbox(slide, 0.34, y + 0.02, SIDEBAR_W - 0.62, item_height - 0.04, vertical_anchor=MSO_ANCHOR.MIDDLE)
        paragraph = box.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        add_text_run(paragraph, chapter, 13.3 if is_active else 11.5, WHITE if is_active else DIM, bold=is_active)


def add_safe_zone(slide) -> None:
    safe_zone = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(SAFE_X),
        Inches(SAFE_Y),
        Inches(SAFE_W),
        Inches(SAFE_H),
    )
    style_shape_fill(safe_zone, WHITE)
    safe_zone.line.color.rgb = LINE
    safe_zone.line.width = Pt(1)


def add_chapter_chip(slide, chapter: str, current_y: float) -> float:
    chip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(SAFE_X + 0.26),
        Inches(current_y),
        Inches(2.3),
        Inches(0.34),
    )
    style_shape_fill(chip, PLACEHOLDER_BG)
    chip.line.fill.background()
    box = add_textbox(slide, SAFE_X + 0.38, current_y + 0.05, 2.0, 0.18, vertical_anchor=MSO_ANCHOR.MIDDLE)
    add_text_run(box.text_frame.paragraphs[0], chapter, 9.5, XMU_BLUE, bold=True)
    return current_y + 0.46


def add_bullet_block(slide, bullet: BulletPoint, index: int, current_y: float, content_bottom: float) -> float:
    x = SAFE_X + 0.28
    w = SAFE_W - 0.56
    keyword_h = estimate_textbox_height(bullet.keyword, w - 0.9, 14.5, padding_in=0.03)
    desc_h = estimate_textbox_height(bullet.description, w - 0.9, 11.4, padding_in=0.03)
    block_h = max(0.72, min(1.12, 0.18 + keyword_h + desc_h))

    if current_y + block_h > content_bottom:
        raise ValueError("Bullet block exceeds content safe zone.")

    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(current_y),
        Inches(w),
        Inches(block_h),
    )
    style_shape_fill(card, CARD_BG)
    card.line.color.rgb = LINE
    card.line.width = Pt(0.9)

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(current_y),
        Inches(0.12),
        Inches(block_h),
    )
    style_shape_fill(accent, XMU_BLUE)
    accent.line.fill.background()

    index_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(x + 0.18),
        Inches(current_y + 0.14),
        Inches(0.28),
        Inches(0.28),
    )
    style_shape_fill(index_box, PLACEHOLDER_BG)
    index_box.line.fill.background()
    idx_text = add_textbox(slide, x + 0.235, current_y + 0.17, 0.18, 0.18, vertical_anchor=MSO_ANCHOR.MIDDLE)
    idx_paragraph = idx_text.text_frame.paragraphs[0]
    idx_paragraph.alignment = PP_ALIGN.CENTER
    add_text_run(idx_paragraph, str(index), 9.0, XMU_BLUE, bold=True)

    keyword_box = add_textbox(slide, x + 0.56, current_y + 0.08, w - 0.72, keyword_h)
    keyword_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    add_text_run(keyword_box.text_frame.paragraphs[0], bullet.keyword, 14.5, XMU_BLUE, bold=True)

    desc_box = add_textbox(slide, x + 0.56, current_y + 0.10 + keyword_h, w - 0.72, desc_h)
    desc_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    add_text_run(desc_box.text_frame.paragraphs[0], bullet.description, 11.4, TEXT)

    return current_y + block_h + 0.14


def add_callout_box(slide, takeaway: str) -> tuple[float, float]:
    callout_h = 0.96
    x = SAFE_X + 0.28
    w = SAFE_W - 0.56
    y = SAFE_Y + SAFE_H - callout_h - 0.22

    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(callout_h),
    )
    style_shape_fill(box, CALLOUT_BG)
    box.line.color.rgb = LINE
    box.line.width = Pt(1)

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(0.12),
        Inches(callout_h),
    )
    style_shape_fill(accent, XMU_BLUE)
    accent.line.fill.background()

    label_box = add_textbox(slide, x + 0.24, y + 0.1, 2.0, 0.18)
    add_text_run(label_box.text_frame.paragraphs[0], "重点结论 | KEY TAKEAWAY", 8.8, DIM, bold=True)

    takeaway_box = add_textbox(slide, x + 0.24, y + 0.32, w - 0.34, 0.42, vertical_anchor=MSO_ANCHOR.MIDDLE)
    takeaway_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    add_text_run(takeaway_box.text_frame.paragraphs[0], takeaway, 17.0, XMU_BLUE, bold=True)
    return y, callout_h


def add_placeholder_card(slide, placeholder_type: str, current_y: float, content_bottom: float) -> float:
    min_h = 1.18
    available_h = content_bottom - current_y
    if placeholder_type == "None" or available_h < 0.52:
        return current_y

    placeholder_h = available_h if available_h >= min_h else available_h
    x = SAFE_X + 0.28
    w = SAFE_W - 0.56
    y = current_y

    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(placeholder_h),
    )
    style_shape_fill(card, PLACEHOLDER_BG)
    card.line.color.rgb = XMU_BLUE
    card.line.width = Pt(1.2)

    title_box = add_textbox(slide, x + 0.18, y + 0.1, 3.0, 0.2)
    add_text_run(title_box.text_frame.paragraphs[0], f"建议配图 | {placeholder_type}", 10.2, XMU_BLUE, bold=True)

    label_box = add_textbox(slide, x + w - 2.7, y + 0.1, 2.5, 0.18)
    label_box.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    add_text_run(label_box.text_frame.paragraphs[0], "Flow-safe visual frame", 8.5, DIM)

    inner_x = x + 0.18
    inner_y = y + 0.36
    inner_w = w - 0.36
    inner_h = placeholder_h - 0.52
    draw_placeholder_visual(slide, placeholder_type, inner_x, inner_y, inner_w, inner_h)
    return y + placeholder_h + 0.12


def draw_label_tile(slide, x: float, y: float, w: float, h: float, title: str, fill: RGBColor) -> None:
    tile = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    style_shape_fill(tile, fill)
    tile.line.color.rgb = LINE
    tile.line.width = Pt(0.8)
    text_box = add_textbox(slide, x + 0.08, y + 0.08, w - 0.16, h - 0.16, vertical_anchor=MSO_ANCHOR.MIDDLE)
    text_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text_run(text_box.text_frame.paragraphs[0], title, 10.5, XMU_BLUE, bold=True)


def draw_arrow_bar(slide, x: float, y: float, w: float, h: float, label: str) -> None:
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    style_shape_fill(arrow, XMU_BLUE)
    arrow.line.fill.background()
    label_box = add_textbox(slide, x, y - 0.16, w, 0.14)
    label_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text_run(label_box.text_frame.paragraphs[0], label, 8.2, DIM)


def draw_placeholder_visual(slide, placeholder_type: str, x: float, y: float, w: float, h: float) -> None:
    if placeholder_type == "ChemDraw":
        draw_label_tile(slide, x + 0.2, y + 0.2, 1.35, h - 0.4, "Ligand", WHITE)
        ring_left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.HEXAGON, Inches(x + 1.78), Inches(y + 0.32), Inches(1.05), Inches(0.84))
        style_shape_fill(ring_left, WHITE)
        ring_left.line.color.rgb = XMU_BLUE
        ring_left.line.width = Pt(1.2)
        ring_right = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.HEXAGON, Inches(x + 3.0), Inches(y + 0.32), Inches(1.05), Inches(0.84))
        style_shape_fill(ring_right, WHITE)
        ring_right.line.color.rgb = XMU_BLUE
        ring_right.line.width = Pt(1.2)
        bridge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 2.68), Inches(y + 0.65), Inches(0.34), Inches(0.04))
        style_shape_fill(bridge, XMU_BLUE)
        bridge.line.fill.background()
        metal_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 4.55), Inches(y + 0.42), Inches(0.7), Inches(0.7))
        style_shape_fill(metal_box, WHITE)
        metal_box.line.color.rgb = XMU_BLUE
        metal_box.line.width = Pt(1.1)
        metal_text = add_textbox(slide, x + 4.68, y + 0.58, 0.42, 0.16, vertical_anchor=MSO_ANCHOR.MIDDLE)
        metal_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        add_text_run(metal_text.text_frame.paragraphs[0], "Fe", 11, XMU_BLUE, bold=True)
        note = add_textbox(slide, x + 5.55, y + 0.28, w - 5.75, h - 0.56)
        note.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        add_text_run(note.text_frame.paragraphs[0], "在此插入 ChemDraw 分子结构图、位点示意或 Reaction Mechanism 路径。", 10.6, TEXT)
        return

    if placeholder_type == "Aspen PFD":
        box_w = 1.12
        draw_label_tile(slide, x + 0.16, y + 0.36, box_w, 0.72, "Feed", WHITE)
        draw_label_tile(slide, x + 1.8, y + 0.36, box_w, 0.72, "R-101", WHITE)
        draw_label_tile(slide, x + 3.44, y + 0.36, box_w, 0.72, "T-201", WHITE)
        draw_label_tile(slide, x + 5.08, y + 0.36, box_w, 0.72, "V-301", WHITE)
        draw_label_tile(slide, x + 6.72, y + 0.36, box_w, 0.72, "PC", WHITE)
        draw_arrow_bar(slide, x + 1.3, y + 0.55, 0.36, 0.18, "F-101")
        draw_arrow_bar(slide, x + 2.94, y + 0.55, 0.36, 0.18, "F-201")
        draw_arrow_bar(slide, x + 4.58, y + 0.55, 0.36, 0.18, "F-301")
        draw_arrow_bar(slide, x + 6.22, y + 0.55, 0.36, 0.18, "F-401")
        note = add_textbox(slide, x + 0.18, y + 1.34, w - 0.36, max(0.4, h - 1.48))
        note.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        add_text_run(note.text_frame.paragraphs[0], "在此插入 Aspen PFD / P&ID 流程图，可继续替换为物流编号、温压条件、回收回路和关键控制阀。", 10.6, TEXT)
        return

    if placeholder_type == "Characterization":
        draw_label_tile(slide, x + 0.16, y + 0.24, 1.65, 0.88, "SEM", WHITE)
        draw_label_tile(slide, x + 1.96, y + 0.24, 1.65, 0.88, "TEM", WHITE)
        draw_label_tile(slide, x + 3.76, y + 0.24, 1.65, 0.88, "XRD", WHITE)
        draw_label_tile(slide, x + 5.56, y + 0.24, 1.65, 0.88, "XPS / FTIR", WHITE)
        note = add_textbox(slide, x + 7.5, y + 0.2, w - 7.68, max(0.48, h - 0.4))
        note.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        add_text_run(note.text_frame.paragraphs[0], "在此插入 SEM/TEM 图、XRD 峰型或 In-situ FTIR 谱图，形成表征证据闭环。", 10.4, TEXT)
        return

    if placeholder_type == "Equipment":
        vessel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 0.36), Inches(y + 0.26), Inches(0.9), Inches(h - 0.52))
        style_shape_fill(vessel, WHITE)
        vessel.line.color.rgb = XMU_BLUE
        vessel.line.width = Pt(1.1)
        exchanger = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 1.72), Inches(y + 0.64), Inches(1.05), Inches(0.56))
        style_shape_fill(exchanger, WHITE)
        exchanger.line.color.rgb = XMU_BLUE
        exchanger.line.width = Pt(1.1)
        tower = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 3.18), Inches(y + 0.18), Inches(0.88), Inches(h - 0.36))
        style_shape_fill(tower, WHITE)
        tower.line.color.rgb = XMU_BLUE
        tower.line.width = Pt(1.1)
        draw_arrow_bar(slide, x + 1.3, y + 0.82, 0.28, 0.16, "S-101")
        draw_arrow_bar(slide, x + 2.84, y + 0.82, 0.24, 0.16, "S-201")
        note = add_textbox(slide, x + 4.45, y + 0.24, w - 4.63, max(0.5, h - 0.38))
        note.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        add_text_run(note.text_frame.paragraphs[0], "在此插入关键设备布置、材质说明、负荷摘要或设备表，突出选型理由与冗余设计。", 10.6, TEXT)
        return

    if placeholder_type == "Risk Matrix":
        cell_w = 0.52
        cell_h = 0.38
        colors = [
            [RGBColor(0xEC, 0xF8, 0xF2), RGBColor(0xEC, 0xF8, 0xF2), RGBColor(0xFF, 0xF6, 0xE8)],
            [RGBColor(0xEC, 0xF8, 0xF2), RGBColor(0xFF, 0xF6, 0xE8), RGBColor(0xFB, 0xEE, 0xEE)],
            [RGBColor(0xFF, 0xF6, 0xE8), RGBColor(0xFB, 0xEE, 0xEE), RGBColor(0xFB, 0xEE, 0xEE)],
        ]
        for row in range(3):
            for col in range(3):
                tile = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                    Inches(x + 0.26 + col * cell_w),
                    Inches(y + 0.28 + row * cell_h),
                    Inches(cell_w - 0.02),
                    Inches(cell_h - 0.02),
                )
                style_shape_fill(tile, colors[row][col])
                tile.line.color.rgb = LINE
        note = add_textbox(slide, x + 2.05, y + 0.22, w - 2.23, max(0.5, h - 0.34))
        note.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        add_text_run(note.text_frame.paragraphs[0], "在此插入 HAZOP / LOPA 结果、事故场景分级、应急响应流程或关键防护屏障。", 10.6, TEXT)
        return

    chart_area = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 0.22), Inches(y + 0.24), Inches(3.0), Inches(h - 0.48))
    style_shape_fill(chart_area, WHITE)
    chart_area.line.color.rgb = LINE
    chart_area.line.width = Pt(0.9)
    axis_x = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 0.5), Inches(y + h - 0.48), Inches(2.28), Inches(0.03))
    style_shape_fill(axis_x, XMU_BLUE)
    axis_x.line.fill.background()
    axis_y = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 0.5), Inches(y + 0.44), Inches(0.03), Inches(h - 0.92))
    style_shape_fill(axis_y, XMU_BLUE)
    axis_y.line.fill.background()
    for idx, height in enumerate((0.38, 0.7, 1.0)):
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 0.82 + idx * 0.48), Inches(y + h - 0.48 - height), Inches(0.24), Inches(height))
        style_shape_fill(bar, (GREEN, ORANGE, XMU_BLUE)[idx])
        bar.line.fill.background()
    note = add_textbox(slide, x + 3.48, y + 0.22, w - 3.66, max(0.5, h - 0.34))
    note.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    add_text_run(note.text_frame.paragraphs[0], "在此插入性能曲线、收率对比、动力学 Profile 或 Techno-Economic Analysis 图表。", 10.6, TEXT)


def render_slide(prs: Presentation, chapters: list[str], spec: SlideSpec) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    add_banner(slide, spec.slide_title)
    add_sidebar(slide, chapters, spec.chapter)
    add_safe_zone(slide)

    callout_y, _ = add_callout_box(slide, spec.key_takeaway)
    content_bottom = callout_y - 0.14
    current_y = add_chapter_chip(slide, spec.chapter, SAFE_Y + 0.2)

    for index, bullet in enumerate(spec.bullet_points, start=1):
        current_y = add_bullet_block(slide, bullet, index, current_y, content_bottom)

    add_placeholder_card(slide, spec.placeholder_type, current_y, content_bottom)


def build_presentation(slides: list[SlideSpec], output_path: str | Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    outline = chapter_outline(slides)
    for spec in slides:
        render_slide(prs, outline, spec)

    output = Path(output_path)
    prs.save(output)
    return output
